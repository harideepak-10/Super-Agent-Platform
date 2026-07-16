"""
CreatePresentationTool — generate a PowerPoint (.pptx) slide deck.

Zone: GREEN — runs automatically, no human approval required.
"""
from __future__ import annotations
import json
import logging
import os
import tempfile
from typing import Any
from core.tools.base_tool import BaseTool, ToolZone

logger = logging.getLogger(__name__)


class CreatePresentationTool(BaseTool):
    """Generate a PowerPoint (.pptx) slide deck from structured content.

    Input::

        {
            "title":    "Q3 Business Review",
            "subtitle": "KRYPSOS — July 2026",
            "author":   "Deepak",
            "slides": [
                {
                    "title":   "Revenue Overview",
                    "content": "Revenue grew 18% YoY to ₹4.2Cr. Operating costs down 12%.",
                    "bullets": ["₹4.2Cr total revenue", "18% YoY growth", "12% cost reduction"]
                },
                {
                    "title":   "Key Highlights",
                    "bullets": ["New client onboarded", "Team expanded to 12", "Product launched"]
                }
            ],
            "theme":    "blue"              # "blue" | "green" | "dark" | "minimal" (default: blue)
        }

    Returns::

        {
            "status":    "created",
            "file_path": "/tmp/Q3_Business_Review.pptx",
            "filename":  "Q3_Business_Review.pptx",
            "slides":    4
        }
    """

    name: str = "create_presentation"
    description: str = (
        "Generate a PowerPoint (.pptx) slide deck. GREEN — runs automatically. "
        "Input JSON: {\"title\": \"Q3 Review\", \"slides\": [{\"title\": \"Revenue\", "
        "\"bullets\": [\"Point 1\", \"Point 2\"]}]}. "
        "Returns file_path — pass to upload_to_drive to save to Google Drive."
    )
    zone: ToolZone = ToolZone.GREEN

    def __init__(self, workspace_id: str | None = None) -> None:
        self._workspace_id = workspace_id

    _THEMES = {
        "blue":    {"bg": (0x1E, 0x3A, 0x5F), "accent": (0x3B, 0x82, 0xF6), "text": (0xFF, 0xFF, 0xFF)},
        "green":   {"bg": (0x06, 0x5F, 0x46), "accent": (0x10, 0xB9, 0x81), "text": (0xFF, 0xFF, 0xFF)},
        "dark":    {"bg": (0x1F, 0x29, 0x37), "accent": (0x8B, 0x5C, 0xF6), "text": (0xFF, 0xFF, 0xFF)},
        "minimal": {"bg": (0xFF, 0xFF, 0xFF), "accent": (0x1E, 0x3A, 0x5F), "text": (0x11, 0x18, 0x27)},
    }

    def run(self, input_str: str) -> str:
        try:
            data = json.loads(input_str) if isinstance(input_str, str) else input_str
        except (json.JSONDecodeError, TypeError):
            return json.dumps({"error": "Invalid input."})

        title    = data.get("title", "Presentation")
        subtitle = data.get("subtitle", "")
        author   = data.get("author", "KRYPSOS AI")
        slides   = data.get("slides", [])
        theme    = data.get("theme", "blue")

        if not slides:
            return json.dumps({"error": "'slides' list is required."})

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs   = Presentation()
            theme_colors = self._THEMES.get(theme, self._THEMES["blue"])
            bg_color     = RGBColor(*theme_colors["bg"])
            accent_color = RGBColor(*theme_colors["accent"])
            text_color   = RGBColor(*theme_colors["text"])

            def set_bg(slide):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = bg_color

            def add_textbox(slide, text, left, top, width, height, font_size=18,
                            bold=False, color=None, align=PP_ALIGN.LEFT):
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf    = txBox.text_frame
                tf.word_wrap = True
                p    = tf.paragraphs[0]
                p.alignment = align
                run  = p.add_run()
                run.text = text
                run.font.size  = Pt(font_size)
                run.font.bold  = bold
                run.font.color.rgb = color or text_color

            # --- Title slide ---
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            set_bg(slide)
            W = prs.slide_width
            H = prs.slide_height
            # Accent bar
            shape = slide.shapes.add_shape(1, 0, int(H * 0.6), W, Pt(4))
            shape.fill.solid()
            shape.fill.fore_color.rgb = accent_color
            shape.line.fill.background()
            add_textbox(slide, title, Inches(0.7), Inches(2), Inches(8.6), Inches(1.5),
                        font_size=36, bold=True, align=PP_ALIGN.CENTER)
            if subtitle:
                add_textbox(slide, subtitle, Inches(0.7), Inches(3.6), Inches(8.6), Inches(0.8),
                            font_size=20, align=PP_ALIGN.CENTER)
            add_textbox(slide, author, Inches(0.7), Inches(4.5), Inches(8.6), Inches(0.5),
                        font_size=14, align=PP_ALIGN.CENTER)

            # --- Content slides ---
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                set_bg(slide)

                # Header bar
                header = slide.shapes.add_shape(1, 0, 0, W, Inches(1.1))
                header.fill.solid()
                header.fill.fore_color.rgb = accent_color
                header.line.fill.background()

                slide_title   = slide_data.get("title", "")
                slide_content = slide_data.get("content", "")
                bullets       = slide_data.get("bullets", [])

                add_textbox(slide, slide_title, Inches(0.4), Inches(0.1), Inches(9.2), Inches(0.9),
                            font_size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

                y = Inches(1.3)
                if slide_content:
                    add_textbox(slide, slide_content, Inches(0.5), y, Inches(9.0), Inches(1.0),
                                font_size=16)
                    y += Inches(1.1)

                for bullet in bullets:
                    add_textbox(slide, f"• {bullet}", Inches(0.7), y, Inches(8.8), Inches(0.5),
                                font_size=15)
                    y += Inches(0.5)

            # Save
            safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)
            filename   = f"{safe_title.replace(' ', '_')}.pptx"
            file_path  = os.path.join(tempfile.gettempdir(), filename)
            prs.save(file_path)

            logger.info("CreatePresentationTool: created %s slides=%d", file_path, len(slides) + 1)
            return json.dumps({
                "status":    "created",
                "file_path": file_path,
                "filename":  filename,
                "slides":    len(slides) + 1,
                "note":      "Pass file_path to upload_to_drive to save to Google Drive.",
            })

        except ImportError:
            return json.dumps({"error": "python-pptx not installed. Run: pip install python-pptx"})
        except Exception as exc:
            logger.exception("CreatePresentationTool failed")
            return json.dumps({"error": str(exc)})

    def to_schema(self) -> dict:
        return {"type": "function", "function": {
            "name": self.name, "description": self.description,
            "parameters": {"type": "object", "properties": {
                "title":    {"type": "string"},
                "subtitle": {"type": "string"},
                "author":   {"type": "string"},
                "theme":    {"type": "string", "enum": ["blue", "green", "dark", "minimal"]},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title":   {"type": "string"},
                            "content": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                        },
                    },
                    "description": "List of slide objects with title and bullets",
                },
            }, "required": ["title", "slides"]},
        }}
